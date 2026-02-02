"""
DATEI: backend/ppt_utils.py
BESCHREIBUNG:
    Enthält Low-Level-Hilfsfunktionen für die Manipulation von PowerPoint-Dateien.
    
    Da die Bibliothek `python-pptx` nicht alle Funktionen nativ unterstützt (z.B. Duplizieren von Folien),
    greifen wir hier teilweise direkt auf die XML-Struktur (OpenXML) zu.
    
    Hauptfunktionen:
    1.  `replace_text_in_shape`: Suchen und Ersetzen von Text in Textfeldern und Tabellen.
    2.  `duplicate_slide`: Erstellt eine exakte Kopie einer Folie inklusive aller Elemente.
    3.  `delete_slide`: Löscht eine Folie aus der Präsentation.
"""

from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy
import re
import time
import random

def replace_text_in_shape(shape, replacements):
    """
    Ersetzt Text in einer Form (Shape) basierend auf einem Dictionary von Ersetzungen.
    
    Argumente:
        shape: Das PowerPoint-Shape-Objekt (Textfeld, Tabelle, etc.).
        replacements: Ein Dictionary der Struktur:
                      { '{{PLATZHALTER}}': { 'text': 'Neuer Wert', 'formatting': {...} } }
    """
    # Früher Abbruch, wenn das Shape keinen Text enthalten kann
    if not shape.has_text_frame and not shape.has_table:
        return

    # Verarbeitung von Tabellen
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text_frame = cell.text_frame
                process_text_frame(text_frame, replacements)
    
    # Verarbeitung von normalen Textfeldern
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, replacements)

def process_text_frame(text_frame, replacements):
    """
    Iteriert durch alle Absätze eines TextFrames und führt Ersetzungen durch.
    """
    for p in text_frame.paragraphs:
        # Optimierung: Wir fassen den Absatz nur an, wenn er Marker ("{{") enthält.
        if "{{" in p.text:
            process_paragraph(p, replacements)

def process_paragraph(p, replacements):
    """
    Kernlogik für das Ersetzen in einem Absatz.
    
    Herausforderung:
    Ein Absatz besteht aus "Runs" (Text-Teilen mit gleicher Formatierung).
    Ein Platzhalter kann über mehrere Runs verteilt sein (z.B. Run1="{{", Run2="Title", Run3="}}").
    Daher bauen wir den Absatz neu auf.
    """
    current_text = p.text
    # Regex-Split, um Platzhalter von statischem Text zu trennen
    # Wir suchen nach Mustern wie {{...}}
    pattern = r"(\{\{.*?\}\})"
    parts = re.split(pattern, current_text)
    
    # Vorprüfung: Haben wir überhaupt eine passende Ersetzung definiert?
    has_match = False
    for part in parts:
        norm_part = " ".join(part.split()) # Leerzeichen normalisieren
        if norm_part in replacements:
            has_match = True
            break
            
    if not has_match:
        return

    # Absatz leeren und neu befüllen
    # p.clear() entfernt alle Runs, behält aber die Absatz-Eigenschaften (Ausrichtung, Abstand etc.) bei.
    p.clear() 
    
    for part in parts:
        if not part: continue
        
        norm_part = " ".join(part.split())
        
        if norm_part in replacements:
            # Es ist ein bekannter Platzhalter -> Ersetzen
            data = replacements[norm_part]
            run = p.add_run()
            run.text = data["text"]
            apply_formatting(run, data.get("formatting", {}))
        else:
            # Es ist statischer Text -> Einfach wieder einfügen
            if part: 
                run = p.add_run()
                # Fix: Vertikale Tabs (\x0b) werden von PPT manchmal als Kästchen (_x000B_) dargestellt.
                # Wir ersetzen sie durch echte Zeilenumbrüche.
                run.text = part.replace("\x0b", "\n")
                
                # Layout-Schutz: Wir setzen eine kleine Schriftgröße (7pt) sicherheitshalber zurück,
                # um zu verhindern, dass Tabellenzeilen durch Formatierungsverlust explodieren.
                run.font.size = Pt(7)

def apply_formatting(run, formatting):
    """
    Wendet Formatierungen (Fett, Größe, Farbe) auf einen Text-Run an.
    """
    if not formatting: return
    font = run.font
    if "bold" in formatting:
        font.bold = formatting["bold"]
    if "font_size" in formatting:
        font.size = Pt(formatting["font_size"])
    if "color" in formatting:
        font.color.rgb = formatting["color"]

def duplicate_slide(prs, source_slide_index):
    """
    Dupliziert die Folie am angegebenen Index und fügt sie am Ende der Präsentation an.
    
    Da python-pptx dies nicht nativ kann, kopieren wir die XML-Elemente.
    """
    source_slide = prs.slides[source_slide_index]
    slide_layout = source_slide.slide_layout
    
    # Neue, leere Folie basierend auf dem gleichen Layout erstellen
    dest_slide = prs.slides.add_slide(slide_layout)
    
    # Alle Formen (Shapes) kopieren
    for shape in source_slide.shapes:
        copy_shape(shape, dest_slide)
        
    return dest_slide

def copy_shape(shape, dest_slide):
    """
    Kopiert eine Form (Shape) auf die Ziel-Folie.
    Nutzt Deep-Copy auf XML-Ebene.
    """
    new_el = copy.deepcopy(shape.element)
    
    # WICHTIG: Jedes Shape muss eine eindeutige ID haben (cNvPr id).
    # Beim bloßen Kopieren hätten wir zwei Shapes mit gleicher ID -> Datei korrupt.
    # Wir generieren daher eine neue, zufällige ID.
    unique_id = int(time.time() * 1000) + random.randint(0, 10000)
    
    # Wir suchen im XML-Baum nach dem Element 'cNvPr' (Non-Visual Properties)
    for desc in new_el.iterdescendants():
        if desc.tag.endswith('cNvPr'):
             # Neue ID setzen
             desc.set('id', str(unique_id))
             # Auch den Namen unique machen (z.B. "Textfeld 12345")
             desc.set('name', desc.get('name') + f" {unique_id}")
             break
    
    # Das neue Element in den XML-Baum der Ziel-Folie einfügen
    dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_el

def delete_slide(prs, index):
    """
    Löscht eine Folie aus der Präsentation anhand ihres Index.
    """
    # Zugriff auf die interne Slide-ID-Liste im XML
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Element entfernen
    xml_slides.remove(slides[index])
