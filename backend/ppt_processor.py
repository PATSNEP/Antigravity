"""
DATEI: backend/ppt_processor.py
BESCHREIBUNG:
    Dies ist das Herzstück der Anwendung.
    Es steuert den gesamten Prozess der PowerPoint-Generierung.
    
    Hauptaufgaben:
    1.  Laden und Filtern der Daten aus der CSV.
    2.  Öffnen der PowerPoint-Vorlage.
    3.  Befüllen der Übersichtsfolie (Slide 1) und der Heatmaps (Slides 2-8).
    4.  Generierung der One-Pager durch Duplizierung von Vorlagenfolien.
    5.  Intelligente Bereinigung ungenutzter Platzhalter.
    6.  Speichern des fertigen Reports.
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re
from datetime import datetime
try:
    from backend.data_loader import load_data
    from backend.ppt_utils import replace_text_in_shape, duplicate_slide, delete_slide
except ImportError:
    from data_loader import load_data
    from ppt_utils import replace_text_in_shape, duplicate_slide, delete_slide

# Konstanten (Entsprechen den Anforderungen des Nutzers)
# Pfad zur Vorlage relativ zum Backend-Ordner
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, "..", "PPTWITHPLACEHOLDERS.pptx")

# Formatierungs-Vorgaben
FMT_TITLE = {"bold": True, "font_size": 7, "color": RGBColor(0, 176, 240)} # Blau
FMT_DATE = {"bold": True, "font_size": 7, "color": RGBColor(0, 0, 0)}     # Schwarz

# Konfiguration der verschiedenen Geschäftsbereiche (Lines of Business)
# Definiert, welche Slides zu welchem Bereich gehören und welche Regex-Muster genutzt werden.
HEATMAP_CONFIGS = [
    {
        "name": "Marketing",
        "filter": "Marketing",
        "slides": [1, 2], # Entspricht Slide 2 & 3 in PowerPoint (0-indiziert)
        "regex_title": r"\{\{Marketing\s+USE\s+CASE\s+Title\s+(\d+)\}\}",
        "fmt_title": "{{{{Marketing USE CASE Title {idx}}}}}",
        "fmt_status": "{{{{StatusupdateUC{idx}Marketing}}}}",
        "key_owner": "{{UseCaseOwnerMarketing}}",
        "fmt_date_d": "{{{{MD{idx}}}}}",
        "fmt_date_a": "{{{{MA{idx}}}}}",
        "fmt_completeness": "{{{{OCM{idx}}}}}",
        "regex_completeness": r"\{\{OCM(\d+)\}\}",
        "regex_date_d": r"\{\{MD(\d+)\}\}",
        "regex_date_a": r"\{\{MA(\d+)\}\}"
    },
    {
        "name": "Sales",
        "filter": "Sales",
        "slides": [3, 4], # Slide 4 & 5
        "regex_title": r"\{\{SALES\s+USE\s+CASE\s+Title\s+(\d+)\}\}",
        "fmt_title": "{{{{SALES USE CASE Title {idx}}}}}",
        "fmt_status": "{{{{StatusupdateUC{idx}Sales}}}}",
        "key_owner": "{{UseCaseOwnerSales}}",
        "fmt_date_d": "{{{{SD{idx}}}}}",
        "fmt_date_a": "{{{{SA{idx}}}}}",
        "fmt_completeness": "{{{{OCS{idx}}}}}",
        "regex_completeness": r"\{\{OCS(\d+)\}\}",
        "regex_date_d": r"\{\{SD(\d+)\}\}",
        "regex_date_a": r"\{\{SA(\d+)\}\}"
    },
    {
        "name": "Compliance",
        "filter": "Compliance",
        "slides": [5], # Slide 6
        "regex_title": r"\{\{Compliance\s+USE\s+CASE\s+Title\s+(\d+)\}\}",
        "fmt_title": "{{{{Compliance USE CASE Title {idx}}}}}",
        "fmt_status": "{{{{StatusupdateUC{idx}Compliance}}}}",
        "key_owner": "{{UseCaseOwnerCompliance}}",
        "fmt_date_d": "{{{{COD{idx}}}}}",
        "fmt_date_a": "{{{{COA{idx}}}}}",
        "fmt_completeness": "{{{{OCC{idx}}}}}",
        "regex_completeness": r"\{\{OCC(\d+)\}\}",
        "regex_date_d": r"\{\{COD(\d+)\}\}",
        "regex_date_a": r"\{\{COA(\d+)\}\}"
    },
    {
        "name": "Customer Success",
        "filter": "Customer Success",
        "slides": [6], # Slide 7
        "regex_title": r"\{\{CS\s+USE\s+CASE\s+Title\s+(\d+)\}\}",
        "fmt_title": "{{{{CS USE CASE Title {idx}}}}}",
        "fmt_status": "{{{{StatusupdateUC{idx}CS}}}}",
        "key_owner": "{{UseCaseOwnerCS}}",
        "fmt_date_d": "{{{{CUD{idx}}}}}",
        "fmt_date_a": "{{{{CUA{idx}}}}}",
        "fmt_completeness": "{{{{OCCS{idx}}}}}",
        "regex_completeness": r"\{\{OCCS(\d+)\}\}",
        "regex_date_d": r"\{\{CUD(\d+)\}\}",
        "regex_date_a": r"\{\{CUA(\d+)\}\}"
    },
    {
        "name": "Finance",
        "filter": "Finance",
        "slides": [7], # Slide 8
        "regex_title": r"\{\{F\s+USE\s+CASE\s+Title\s+(\d+)\}\}",
        "fmt_title": "{{{{F USE CASE Title {idx}}}}}",
        "fmt_status": "{{{{StatusupdateUC{idx}F}}}}",
        "key_owner": "{{UseCaseOwnerF}}",
        "fmt_date_d": "{{{{FD{idx}}}}}",
        "fmt_date_a": "{{{{FA{idx}}}}}",
        "fmt_completeness": "{{{{OCF{idx}}}}}",
        "regex_completeness": r"\{\{OCF(\d+)\}\}",
        "regex_date_d": r"\{\{FD(\d+)\}\}",
        "regex_date_a": r"\{\{FA(\d+)\}\}"
    }
]

def process_ppt(csv_path, output_folder):
    """
    Hauptfunktion: Verarbeitet die PowerPoint mit den Daten aus der CSV.
    
    Argumente:
        csv_path: Pfad zur hochgeladenen CSV-Datei.
        output_folder: Pfad, wo der fertige Report gespeichert werden soll.
        
    Rückgabe:
        Dateiname des generierten Reports.
    """
    print(f"Verarbeite CSV: {csv_path}")
    
    # 1. Daten laden
    # data_loader gruppiert die Daten nach 'line_of_business' oder 'business_unit'
    raw_data = load_data(csv_path) 
    
    # Initialisierung der Ersetzungen für Slide 1
    replacements = {}
    
    # Konfiguration für die Übersicht (Slide 1)
    # Mapping von LoB-Namen zu den spezifischen Platzhaltern auf Slide 1
    LOB_CONFIGS = [
        {
            "name": "Marketing",
            "filter": "Marketing",
            "p_title": "Marketing USE CASE Title {}",
            "p_del": "MD{}",
            "p_adopt": "MA{}"
        },
        {
            "name": "Sales",
            "filter": "Sales",
            "p_title": "SALES USE CASE Title {}", # Achtung: Großschreibung im Template
            "p_del": "SD{}",
            "p_adopt": "SA{}"
        },
        {
            "name": "Compliance",
            "filter": "Compliance",
            "p_title": "Compliance USE CASE Title {}",
            "p_del": "COD{}",
            "p_adopt": "COA{}"
        },
        {
            "name": "Customer Success",
            "filter": "Customer Success",
            "p_title": "Customer Success USE CASE Title {}",
            "p_del": "CUD{}",
            "p_adopt": "CUA{}"
        },
        {
            "name": "Finance",
            "filter": "Finance",
            "p_title": "Finance USE CASE Title {}",
            "p_del": "FD{}",
            "p_adopt": "FA{}"
        }
    ]
    
    # Flache Liste aller Cases erstellen, um später einfacher zu filtern
    all_cases = []
    for cases in raw_data.values():
        all_cases.extend(cases)
        
    for config in LOB_CONFIGS:
        # 1. Grober Filter nach Business Unit
        lob_cases = [c for c in all_cases if config["filter"] in getattr(c, "business_unit", "")]
        
        # 2. Strikter Filter für Slide 1 (Anforderung: Nur "CDP Business Adoption" anzeigen)
        # "CDP Foundational Use Cases" werden hier ignoriert.
        slide1_display_cases = [
            c for c in lob_cases 
            if getattr(c, "use_case_type", "").strip() == "CDP Business Adoption"
        ]
        
        print(f"LoB: {config['name']} | Gefunden: {len(lob_cases)} | Anzeige (Business Adoption): {len(slide1_display_cases)}")
        
        for i, case in enumerate(slide1_display_cases):
            # i+1, da Platzhalter bei 1 beginnen
            idx = i + 1
            
            # Mapping der Attribute zu Platzhaltern
            # Titel
            key_title = "{{" + config["p_title"].format(idx) + "}}"
            replacements[key_title] = {
                "text": case.title,
                "formatting": FMT_TITLE
            }
            # Lieferdatum
            key_del = "{{" + config["p_del"].format(idx) + "}}"
            replacements[key_del] = {"text": case.delivery_date, "formatting": FMT_DATE}
            
            # Adoptionsdatum
            key_adopt = "{{" + config["p_adopt"].format(idx) + "}}"
            replacements[key_adopt] = {"text": case.adoption_date, "formatting": FMT_DATE}
    
    # 2. Vorlage öffnen
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(f"Vorlage nicht gefunden unter: {TEMPLATE_PATH}")
        
    prs = Presentation(TEMPLATE_PATH)

    # 3. Slide 1 Logik (Übersicht)
    # Anwenden der generischen Ersetzungen auf Slide 1 (Index 0).
    print("Verarbeite Slide 1 (Übersicht)...")
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        replace_text_in_shape(shape, replacements)
        
    # 4. Slide 2-8 Logik (Heatmaps / Status Slides)
    # Iteriert durch die konfigurierten Heatmap-Bereiche

    for config in HEATMAP_CONFIGS:
        # Filter: Nur Business Adoption Cases der jeweiligen LoB
        cases = [
            c for c in all_cases 
            if config["filter"] in getattr(c, "business_unit", "") 
            and getattr(c, "use_case_type", "").strip() == "CDP Business Adoption"
        ]
        
        print(f"Verarbeite Heatmaps für {config['name']} ({len(cases)} Fälle gefunden)...")
        
        for slide_idx in config["slides"]:
            if slide_idx >= len(prs.slides): continue
            
            slide = prs.slides[slide_idx]
            
            # Durchlaufe Formen auf der Folie
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        
                        # 4.1 Identifiziere den Case für diese Zeile
                        # Wir suchen nach dem Titel-Platzhalter (z.B. {{Sales USE CASE Title 1}})
                        # um zu wissen, welche ID (1, 2, 3...) diese Zeile repräsentiert.
                        row_case_idx = -1
                        row_case = None
                        
                        # Scan in Spalte 0 nach dem Titel
                        if len(row.cells) > 0:
                            c0_text = row.cells[0].text_frame.text
                            match = re.search(config["regex_title"], c0_text, re.IGNORECASE)
                            if match:
                                idx_found = int(match.group(1))
                                row_case_idx = idx_found - 1 # 0-basiert
                                if 0 <= row_case_idx < len(cases):
                                    row_case = cases[row_case_idx]

                        # 4.2 Heatmap Einfärbung (Traffic Lights)
                        if row_case:
                            # Parse Status-Schritt (z.B. "7. Technical GoLive") -> Schritt 7
                            hm_status_str = getattr(row_case, "heatmap_status", "").strip()
                            current_step = 0
                            step_match = re.match(r"^(\d+)\.", hm_status_str)
                            if step_match:
                                current_step = int(step_match.group(1))
                            
                            # Farben definieren
                            COLOR_LIGHT_GREEN = RGBColor(226, 239, 217) # Erledigt
                            COLOR_DARK_GREEN = RGBColor(87, 162, 55)    # Aktuell
                            COLOR_WHITE = RGBColor(255, 255, 255)       # Offen
                            
                            # Iteriere Heatmap-Spalten (1 bis 8)
                            for step_col in range(1, 9):
                                if step_col >= len(row.cells): break
                                
                                cell = row.cells[step_col]
                                cell.fill.solid()
                                
                                if step_col < current_step:
                                    cell.fill.fore_color.rgb = COLOR_LIGHT_GREEN
                                elif step_col == current_step:
                                    cell.fill.fore_color.rgb = COLOR_DARK_GREEN
                                else:
                                    cell.fill.fore_color.rgb = COLOR_WHITE
                        
                        # 4.3 Text-Ersetzungen durchführen
                        for cell in row.cells:
                            process_heatmap_cell(cell.text_frame, cases, config)
                            process_completeness_placeholder(cell.text_frame, cases, config)
                            process_date_placeholders(cell.text_frame, cases, config)
                
                # Auch Textfelder außerhalb von Tabellen verarbeiten
                if shape.has_text_frame:
                    process_heatmap_cell(shape.text_frame, cases, config)
                    process_completeness_placeholder(shape.text_frame, cases, config)
                    process_date_placeholders(shape.text_frame, cases, config)
    
    # 5. Slide 9 & 10 Logik (Foundational Use Cases)
    # Filter: Type="CDP Foundational Use Case" (Unabhängig von Business Unit)
    foundational_cases = [
        c for c in all_cases 
        if getattr(c, "use_case_type", "").strip() == "CDP Foundational Use Case"
    ]
    print(f"Verarbeite Foundational Cases ({len(foundational_cases)} Fälle gefunden)...")
    
    # Slides für Foundational Cases (Indices hängen von Finance LoB ab)
    foundational_slides = [8, 9] 
    
    # Statistik für Overview Message berechnen
    total_foundational = len(foundational_cases)
    positive_count = 0
    for c in foundational_cases:
        status_val = getattr(c, "traffic_light", "").strip().lower()
        # Grün oder Grau (oder leer) gilt als "on track"
        if "green" in status_val or "grey" in status_val or "gray" in status_val or status_val == "":
            positive_count += 1
            
    overview_msg = ""
    if total_foundational > 0:
        percent_positive = (positive_count / total_foundational) * 100
        if percent_positive >= 80:
            overview_msg = "All CDP Foundational Use Cases are on track and will enable business adoption."
        else:
            overview_msg = f"Only {int(percent_positive)}% CDP Foundational Use Cases are on track and will enable business adoption."
    else:
        overview_msg = "No Foundational Use Cases found."

    # Spezifische Ersetzungen für Foundational Cases
    f_replacements = {}
    
    # Overview Messages
    f_replacements["{{AIOverviewMessage1}}"] = {"text": overview_msg, "formatting": {"bold": False, "font_size": 11, "color": RGBColor(0,0,0)}}
    f_replacements["{{AIOverviewMessage2}}"] = {"text": overview_msg, "formatting": {"bold": False, "font_size": 11, "color": RGBColor(0,0,0)}}
    
    for i, case in enumerate(foundational_cases):
        idx = i + 1 # 1-basierter Index
        
        # Titel
        f_replacements[f"{{{{Foundational Use Case Title {idx}}}}}"] = {
            "text": case.title,
            "formatting": FMT_TITLE
        }
        
        # Owner
        f_replacements[f"{{{{Foundational Use Case Owner {idx}}}}}"] = {
            "text": case.owner,
            "formatting": {"font_size": 7, "color": RGBColor(0,0,0)}
        }
        
        # Overall Status
        f_replacements[f"{{{{Overall Status FUC {idx}}}}}"] = {
            "text": getattr(case, "overall_status", "N/A"),
            "formatting": {"font_size": 7, "color": RGBColor(0,0,0)}
        }
        
        
    for slide_idx in foundational_slides:
        if slide_idx >= len(prs.slides): continue
        slide = prs.slides[slide_idx]
        
        for shape in slide.shapes:
            # Standard Text-Ersetzung
            replace_text_in_shape(shape, f_replacements)
            
            # Ampel-Färbung (Traffic Light) via {{prX}} Platzhalter
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_traffic_light_placeholder(cell, foundational_cases)
            
            if shape.has_text_frame:
                process_traffic_light_placeholder(shape, foundational_cases)

    
    # 6. One-Pager Generierung
    # Hier werden Folien dynamisch dupliziert und befüllt.
    
    FMT_OP_TEXT = {"font_size": 10, "color": RGBColor(0,0,0)}
    
    # Sortiere Fälle für One-Pager
    ordered_cases = []
    
    for config in HEATMAP_CONFIGS:
        lob_cases = []
        for c in all_cases:
            if config["filter"] in getattr(c, "business_unit", ""):
                lob_cases.append(c)
        ordered_cases.extend(lob_cases)

    # Start-Index für One-Pager (Slide 11 ist Index 10)
    start_op_index = 10
    print(f"Generiere One-Pagers für {len(ordered_cases)} Fälle (Start auf Slide {start_op_index+1})...")
    
    cases_processed = 0
    
    for i, target_uc in enumerate(ordered_cases):
        slide_idx = start_op_index + i
        
        # Prüfen, ob noch genug Vorlagen-Folien da sind (oder dynamisch erzeugen)
        if slide_idx >= len(prs.slides):
            print(f"WARNUNG: Nicht genug Folien für One-Pager! Stoppe bei Fall {i+1}.")
            break
            
        slide = prs.slides[slide_idx]
        
        # One-Pager Platzhalter (statisch im Template)
        op_replacements = {
            "{{UseCaseOnePagerTitel1}}": {"text": target_uc.title, "formatting": {"bold": True, "color": RGBColor(0, 176, 240)}},
            "{{UseCaseOnePagerPB1}}": {"text": target_uc.problem_statement, "formatting": FMT_OP_TEXT},
            "{{UseCaseOnePagerScope1}}": {"text": target_uc.scope, "formatting": FMT_OP_TEXT},
            "{{UseCaseOnePagerV&KPI1}}": {"text": target_uc.value_kpis, "formatting": FMT_OP_TEXT},
            "{{UseCaseOnePagerBU1}}": {"text": target_uc.line_of_business, "formatting": FMT_OP_TEXT}, 
            "{{UseCaseOnePagerBSU1}}": {"text": target_uc.business_unit, "formatting": FMT_OP_TEXT}, 
            "{{UseCaseOnePagerOwner1}}": {"text": target_uc.owner, "formatting": FMT_OP_TEXT},
            "{{UseCaseOnePagerScopeBC}}": {"text": target_uc.business_contacts, "formatting": FMT_OP_TEXT},
            "{{UseCaseOnePagerScopeAFK}}": {"text": target_uc.affected_key_users, "formatting": FMT_OP_TEXT},
        }
        
        # Folie befüllen
        for shape in slide.shapes:
            replace_text_in_shape(shape, op_replacements)
            
        cases_processed += 1

    # 7. Ungenutzte Slides löschen
    # Wenn wir weniger Fälle als Vorlagen-Slides haben, entfernen wir den Rest.
    
    last_filled_index = start_op_index + cases_processed - 1
    total_slides = len(prs.slides)
    
    delete_start = last_filled_index + 1
    
    if delete_start < total_slides:
        print(f"Entferne ungenutzte Slides von Index {delete_start} bis {total_slides-1}...")
        # Rückwärts löschen, um Index-Verschiebungen zu vermeiden
        for idx in range(total_slides - 1, delete_start - 1, -1):
            delete_slide(prs, idx)
    
    print(f"One-Pager Generierung abgeschlossen. {cases_processed} Folien befüllt.")

    # 8. Auto-Cleanup: Entferne ALLE verbliebenen Platzhalter {{...}}
    # Dies ist wichtig für ein sauberes Endprodukt.
    cleanup_unused_placeholders(prs)

    # 9. Speichern des Outputs
    # Dateiname mit Zeitstempel (Date + Time), um Caching-Probleme zu verhindern.
    # Format: CDP_USECASE_AUTOREPORT_JJJJ-MM-TT_HH-MM-SS.pptx
    timestamp_str = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    output_filename = f"CDP_USECASE_AUTOREPORT_{timestamp_str}.pptx"
    output_path = os.path.join(output_folder, output_filename)
    prs.save(output_path)
    
    return output_filename

def process_heatmap_cell(text_frame, cases, config):
    """
    Hilfsfunktion: Scannt ein Textfeld nach Titeln und führt kontextuelle Ersetzungen durch.
    """
    try:
        from backend.ppt_utils import process_text_frame
    except ImportError:
        from ppt_utils import process_text_frame
    
    text = text_frame.text
    # Suche nach Titel-Platzhalter (z.B. {{Marketing USE CASE Title 1}})
    match = re.search(config["regex_title"], text, re.IGNORECASE)
    
    if match:
        idx = int(match.group(1))
        case_idx = idx - 1
        
        if 0 <= case_idx < len(cases):
            case = cases[case_idx]
            
            replacements = {}
            
            # Titel, Status, Owner und Daten vorbereiten
            key_title = config["fmt_title"].format(idx=idx)
            replacements[key_title] = {"text": case.title, "formatting": FMT_TITLE}
            
            key_status = config["fmt_status"].format(idx=idx)
            replacements[key_status] = {"text": getattr(case, "status_update", "N/A"), "formatting": {"font_size": 7, "color": RGBColor(0,0,0)}}
            
            replacements[config["key_owner"]] = {"text": case.owner, "formatting": {"font_size": 7, "color": RGBColor(0,0,0)}}
            
            # Datum
            FMT_HM_DATE = {"font_size": 10, "bold": False, "color": RGBColor(0,0,0)} # Anforderung: Nicht fett, Größe 10
            key_date_d = config["fmt_date_d"].format(idx=idx)
            replacements[key_date_d] = {"text": case.delivery_date, "formatting": FMT_HM_DATE}
            
            key_date_a = config["fmt_date_a"].format(idx=idx)
            replacements[key_date_a] = {"text": case.adoption_date, "formatting": FMT_HM_DATE}
            
            # Completeness
            if "fmt_completeness" in config:
                key_comp = config["fmt_completeness"].format(idx=idx)
                comp_val = getattr(case, "overall_completeness", "")
                comp_fmt = {"font_size": 10, "color": RGBColor(0,0,0), "bold": False}
                if "100%" in str(comp_val):
                    comp_fmt["color"] = RGBColor(87, 162, 55) # Grün bei 100%
                replacements[key_comp] = {"text": comp_val, "formatting": comp_fmt}
            
            process_text_frame(text_frame, replacements)

def process_traffic_light_placeholder(shape_or_cell, cases):
    """
    Prüft auf Ampel-Platzhalter {{prX}} und färbt den Hintergrund entsprechend ein.
    Der Text des Platzhalters wird danach entfernt.
    """
    if not hasattr(shape_or_cell, "text_frame"): return
    
    text = shape_or_cell.text_frame.text
    match = re.search(r"\{\{pr(\d+)\}\}", text, re.IGNORECASE)
    
    if match:
        idx = int(match.group(1))
        c_idx = idx - 1
        
        if 0 <= c_idx < len(cases):
            case = cases[c_idx]
            color_val = getattr(case, "traffic_light", "").strip().lower()
            
            final_color = RGBColor(200, 200, 200) # Default Grau
            
            if "green" in color_val:
                final_color = RGBColor(87, 162, 55)
            elif "red" in color_val:
                final_color = RGBColor(255, 0, 0)
            elif "yellow" in color_val:
                final_color = RGBColor(247, 203, 84)
            elif "grey" in color_val or "gray" in color_val:
                final_color = RGBColor(128, 128, 128)
            
            # Farbe anwenden
            shape_or_cell.fill.solid()
            shape_or_cell.fill.fore_color.rgb = final_color
            
            # Text entfernen
            shape_or_cell.text_frame.text = ""

def process_completeness_placeholder(text_frame, cases, config):
    """
    Verarbeitet Completeness-Platzhalter (z.B. {{OCM1}}) unabhängig vom Titel-Kontext.
    """
    if "regex_completeness" not in config: return
    try:
        from backend.ppt_utils import process_text_frame
    except ImportError:
        from ppt_utils import process_text_frame
        
    text = text_frame.text
    match = re.search(config["regex_completeness"], text, re.IGNORECASE)
    if match:
        idx = int(match.group(1))
        case_idx = idx - 1
        if 0 <= case_idx < len(cases):
            case = cases[case_idx]
            key_comp = config["fmt_completeness"].format(idx=idx)
            comp_val = getattr(case, "overall_completeness", "")
            comp_fmt = {"font_size": 10, "color": RGBColor(0,0,0), "bold": False}
            if "100%" in str(comp_val):
                comp_fmt["color"] = RGBColor(87, 162, 55)
            process_text_frame(text_frame, {key_comp: {"text": comp_val, "formatting": comp_fmt}})

def process_date_placeholders(text_frame, cases, config):
    """
    Verarbeitet Datums-Platzhalter (z.B. {{MD1}}) unabhängig vom Titel-Kontext.
    """
    if "regex_date_d" not in config or "regex_date_a" not in config: return
    try:
        from backend.ppt_utils import process_text_frame
    except ImportError:
        from ppt_utils import process_text_frame
        
    text = text_frame.text
    replacements = {}
    found_any = False
    FMT_HM_DATE = {"font_size": 10, "bold": False, "color": RGBColor(0,0,0)}

    match_d = re.search(config["regex_date_d"], text, re.IGNORECASE)
    if match_d:
        idx = int(match_d.group(1))
        case = cases[idx - 1] if 0 <= (idx - 1) < len(cases) else None
        if case:
            key = config["fmt_date_d"].format(idx=idx)
            replacements[key] = {"text": case.delivery_date, "formatting": FMT_HM_DATE}
            found_any = True
            
    match_a = re.search(config["regex_date_a"], text, re.IGNORECASE)
    if match_a:
        idx = int(match_a.group(1))
        case = cases[idx - 1] if 0 <= (idx - 1) < len(cases) else None
        if case:
            key = config["fmt_date_a"].format(idx=idx)
            replacements[key] = {"text": case.adoption_date, "formatting": FMT_HM_DATE}
            found_any = True
            
    if found_any:
        process_text_frame(text_frame, replacements)

def cleanup_unused_placeholders(prs):
    """
    Iteriert durch alle Folien und Formen und entfernt verbliebene Platzhalter {{...}}.
    Nutzt eine Layout-sichere Methode ("Smart Run Clearing").
    """
    pattern = re.compile(r"\{\{.*?\}\}", re.DOTALL)
    print("Führe Cleanup durch: Entferne ungenutzte Platzhalter...")
    cleaned_count = 0
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    def iter_shapes(shapes):
        """Rekursiver Iterator für Gruppierte Formen."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(shape.shapes)
            else:
                yield shape
    
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            def clean_frame(tf):
                count = 0
                for p in tf.paragraphs:
                    stripped_text = p.text.strip()
                    # Strategie 1: Absatz enthält NUR einen Platzhalter
                    if stripped_text.startswith("{{") and stripped_text.endswith("}}"):
                         if pattern.fullmatch(stripped_text) or pattern.search(stripped_text):
                             # Sicheres Leeren: Ersetze ersten Run durch Leerzeichen (behält Format), leere den Rest
                             if len(p.runs) > 0:
                                 p.runs[0].text = " "
                                 for r in p.runs[1:]:
                                     r.text = ""
                                 count += 1
                                 continue
                    
                    # Strategie 2: Gemischter Inhalt (Fallback, nur wenn sicher)
                    for run in p.runs:
                        if "{{" in run.text:
                             new_text, n = pattern.subn(" ", run.text)
                             if n > 0:
                                 run.text = new_text
                                 count += n
                return count

            if shape.has_text_frame:
                if shape.text_frame.text:
                    cleaned_count += clean_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame.text:
                            cleaned_count += clean_frame(cell.text_frame)
                                
    print(f"Cleanup abgeschlossen. {cleaned_count} Platzhalter-Fragmente entfernt.")
