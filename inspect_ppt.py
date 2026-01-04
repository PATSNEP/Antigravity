from pptx import Presentation
import os

filename = "Final_Report.pptx"
cwd = "/Users/patrickschnepf/Desktop/Master WINF/1 Semester/Projekt DT/Antigravity"
path = os.path.join(cwd, filename)

prs = Presentation(path)

def print_shape_text(shape, indent=0):
    indent_str = " " * indent
    if hasattr(shape, "text") and shape.text.strip():
        print(f"{indent_str}Shape: {shape.name}, Text: {shape.text}")
        if "{{" in shape.text:
             print(f"{indent_str}  FOUND PLACEHOLDER IN: {shape.text}")
    
    if shape.shape_type == 6:  # Group
        print(f"{indent_str}Group: {shape.name}")
        for child in shape.shapes:
            print_shape_text(child, indent + 2)
            
    if shape.shape_type == 19: # Table
        print(f"{indent_str}Table: {shape.name}")
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame.text.strip():
                     print(f"{indent_str}  Cell Text: {cell.text_frame.text}")
                     if "{{" in cell.text_frame.text:
                            print(f"{indent_str}    FOUND PLACEHOLDER IN: {cell.text_frame.text}")

for i, slide in enumerate(prs.slides):
    print(f"Slide {i+1}:")
    for shape in slide.shapes:
        print_shape_text(shape, 2)
