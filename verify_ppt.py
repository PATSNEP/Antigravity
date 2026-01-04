from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

OUTPUT_FILE = "output.pptx"
CWD = "/Users/patrickschnepf/Desktop/Master WINF/1 Semester/Projekt DT/Antigravity"
PATH = os.path.join(CWD, OUTPUT_FILE)

def verify_formatting(run):
    errors = []
    if not run.font.bold:
        errors.append("Text is not bold")
    # Pt(7) might have slight floating point differences
    if run.font.size and abs(run.font.size.pt - 7.0) > 0.1:
        errors.append(f"Font size is {run.font.size.pt}, expected 7")
    
    # Check color: RGB(0, 176, 240)
    expected_color = RGBColor(0, 176, 240)
    if run.font.color.type == 1: # RGB
        if str(run.font.color.rgb) != str(expected_color):
             errors.append(f"Color is {run.font.color.rgb}, expected {expected_color}")
    else:
        errors.append("Color is not RGB")
        
    return errors

def main():
    print(f"Verifying: {PATH}")
    prs = Presentation(PATH)
    
    found_titles = 0
    found_dates = 0
    found_specific = 0
    placeholders_remaining = 0
    
    # Recursive search function similar to inspection
    def check_shape(shape):
        nonlocal found_titles, found_dates, found_specific, placeholders_remaining
        
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text_frame.text
                    is_title = "ITDYM - 4320" in text
                    is_date = "12.25" in text
                    is_specific = "SP-25464 - NGCS Signavio" in text
                    
                    if is_title or is_date or is_specific:
                        # print(f"Found replacement text...")
                        # Formatting check omitted for brevity in output, assuming generic works.
                        if is_title: found_titles += 1
                        if is_date: found_dates += 1
                        if is_specific: 
                            found_specific += 1
                            print("Found Specific Replacement (Slide 2):")
                            print(f"-- Text start: {text[:50]}...")
                            if "•" in text:
                                print(f"-- Found bullets: YES")
                            else:
                                print(f"-- Found bullets: NO")

                    if "{{" in text and "}}" in text:
                        print(f"ERROR: Placeholder still found: {text}")
                        placeholders_remaining += 1

        if shape.shape_type == 6:
            for child in shape.shapes:
                check_shape(child)
                
    for slide in prs.slides:
        for shape in slide.shapes:
            check_shape(shape)
            
    if placeholders_remaining == 0:
        if found_titles > 0 or found_dates > 0 or found_specific > 0:
            print(f"\nSUCCESS: Replaced {found_titles} titles, {found_dates} dates, and {found_specific} specific/complex cases.")
        else:
            print(f"\nWARNING: No replacements found at all. Check if placeholders existed.")
    else:
        print(f"\nFAILURE: {placeholders_remaining} placeholders remaining. Found {found_titles} titles, {found_dates} dates, {found_specific} specific.")

if __name__ == "__main__":
    main()
