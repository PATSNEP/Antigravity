from pptx import Presentation
from pptx.util import Inches

def create_mock_template():
    prs = Presentation()
    
    # --- Slide 1: Overview (Generic LoB Placeholders) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    title = slide.shapes.title
    title.text = "Overview Slide"
    
    # Create Table
    rows, cols = 5, 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
    
    # Fill generic placeholders for Marketing (just a sample)
    # Row 0: Marketing USE CASE Title 1, {{MD1}}, {{MA1}}
    table.cell(0, 0).text = "{{Marketing USE CASE Title 1}}"
    table.cell(0, 1).text = "{{MD1}}"
    table.cell(0, 2).text = "{{MA1}}"
    
    table.cell(1, 0).text = "{{Marketing USE CASE Title 2}}"
    table.cell(1, 1).text = "{{MD2}}"
    table.cell(1, 2).text = "{{MA2}}"
    
    # Sales
    table.cell(2, 0).text = "{{SALES USE CASE Title 1}}"
    table.cell(2, 1).text = "{{SD1}}"
    table.cell(2, 2).text = "{{SA1}}"
    
    # --- Slide 2: Heatmap (Complex Placeholders) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Heatmap Slide"
    
    # Table for Heatmap
    table2 = slide2.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(3)).table
    
    # Complex Cell: Marketing UC 1
    # Cell content: {{Marketing USE CASE Title 1}}\n{{UseCaseOwnerMarketing}}\n{{StatusupdateUC1Marketing}}
    cell = table2.cell(0, 0)
    cell.text = "{{Marketing USE CASE Title 1}}\n{{UseCaseOwnerMarketing}}\n{{StatusupdateUC1Marketing}}"
    
    # Complex Cell: Marketing UC 2
    cell = table2.cell(1, 0)
    cell.text = "{{Marketing USE CASE Title 2}}\n{{UseCaseOwnerMarketing}}\n{{StatusupdateUC2Marketing}}"
    
    # --- Slide 3: One Pager Template ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1]) # Bullet layout
    slide3.shapes.title.text = "{{UseCaseOnePagerTitel1}}"
    
    # Add text box for body
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.text = "Problem: {{UseCaseOnePagerPB1}}"
    p = tf.add_paragraph()
    p.text = "Scope: {{UseCaseOnePagerScope1}}"
    p = tf.add_paragraph()
    p.text = "Value: {{UseCaseOnePagerV&KPI1}}"
    p = tf.add_paragraph()
    p.text = "Owner: {{UseCaseOnePagerOwner1}}"
    
    # Save
    prs.save("mock_template.pptx")
    print("Created mock_template.pptx")

if __name__ == "__main__":
    create_mock_template()
