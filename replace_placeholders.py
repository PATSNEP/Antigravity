from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import copy

import re

# Configuration
INPUT_FILE = "SAP_IDM_Customer Data Management_CDP_Adoption Reporting_October2025 (2).pptx"
OUTPUT_FILE = "output.pptx"
CWD = "/Users/patrickschnepf/Desktop/Master WINF/1 Semester/Projekt DT/Antigravity"

# Replacement Configuration
TITLE_TEXT = "ITDYM - 4320"
DATE_TEXT = "12.25"

# Regex for Date Placeholders: MD#, MA#, SD#, SA#, CUD#, CUA#, COD#, COA#
# Matches M, S, CU, CO followed by D or A, followed by digits
DATE_REGEX = r"\{\{(M|S|CU|CO)(D|A)\d+\}\}"

REPLACEMENT_FORMATTING = {
    "bold": True,
    "font_size": 7,
    "color": RGBColor(0, 176, 240)
}

def apply_formatting(run, formatting):
    """Applies formatting to a text run."""
    font = run.font
    if "bold" in formatting:
        font.bold = formatting["bold"]
    if "font_size" in formatting:
        font.size = Pt(formatting["font_size"])
    if "color" in formatting:
        font.color.rgb = formatting["color"]

def process_shape(shape):
    """Recursively process shapes to find and replace text."""
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text_frame = cell.text_frame
                current_text = text_frame.text.strip() # Don't strip too much if newlines matter, but here we rebuild
                
                # Check for Slide 2 specific combined placeholders
                if "{{Marketing USE CASE Title 4}}" in current_text and "{{UseCaseOwnerMarketing}}" in current_text:
                    print(f"Replacing Slide 2 specific cell...")
                    text_frame.clear()
                    
                    # 1. Marketing USE CASE Title 4
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "SP-25464 - NGCS Signavio"
                    apply_formatting(run, REPLACEMENT_FORMATTING)
                    
                    # 2. UseCaseOwnerMarketing
                    p = text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = "Niting Singh"
                    run.font.size = Pt(7)
                    # p.space_before = Pt(6) # Optional spacing
                    
                    # 3. StatusupdateUC4Marketing
                    status_text = "MQL replication from HubSpot to CDP recently completed; now in IT testing Lead replication from CDP to HubSpot completed; now in IT testing Consent replication completed Business E2E testing planned in September/October"
                    parts = status_text.split(";")
                    for part in parts:
                        p = text_frame.add_paragraph()
                        p.level = 0 # Indent level
                        # There isn't a simple "bullet" property on p in python-pptx effectively without XML manipulation sometimes,
                        # but we can try letting PPT handle it or using a character. 
                        # However, p.level implies a list level. To force a bullet, we might need a distinct approach.
                        # Simple approach: Bullet character manually if p.level doesn't trigger auto-bullets (dependent on master slide).
                        # Let's try adding a bullet char manually for safety if master doesn't enforce it.
                        run = p.add_run()
                        run.text = "• " + part.strip()
                        run.font.size = Pt(7)
                        
                    continue # Skip normal processing for this cell

                # Check for regex match: {{...}}
                match = re.search(r"\{\{.*?\}\}", current_text)
                if match:
                    print(f"Replacing placeholder in table cell: '{current_text}'")
                    # Clear existing content and replace
                    text_frame.clear() 
                    
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    
                    # Determine replacement text
                    if re.match(DATE_REGEX, current_text):
                        run.text = DATE_TEXT
                    else:
                        run.text = TITLE_TEXT
                        
                    apply_formatting(run, REPLACEMENT_FORMATTING)

    if shape.shape_type == 6:  # Group
        for child in shape.shapes:
            process_shape(child)
            
    if hasattr(shape, "text_frame") and shape.text_frame:
        text_frame = shape.text_frame
        current_text = text_frame.text.strip()
         
        match = re.search(r"\{\{.*?\}\}", current_text)
        if match:
             print(f"Replacing placeholder in text box: '{current_text}'")
             text_frame.clear()
             p = text_frame.paragraphs[0]
             run = p.add_run()
             
             if re.match(DATE_REGEX, current_text):
                run.text = DATE_TEXT
             else:
                run.text = TITLE_TEXT
                
             apply_formatting(run, REPLACEMENT_FORMATTING)

def main():
    input_path = os.path.join(CWD, INPUT_FILE)
    output_path = os.path.join(CWD, OUTPUT_FILE)
    
    print(f"Opening presentation: {input_path}")
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape)
            
    print(f"Saving presentation to: {output_path}")
    prs.save(output_path)
    print("Done.")

if __name__ == "__main__":
    main()
