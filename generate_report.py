import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from data_loader import load_data
from ppt_utils import replace_text_in_shape, duplicate_slide, process_text_frame

# Config
INPUT_PPTX = "mock_template.pptx"
OUTPUT_PPTX = "Final_Report.pptx"
CSV_FILE = "mock_data.csv"

# Formatting Constants
FMT_TITLE = {"bold": True, "font_size": 7, "color": RGBColor(0, 176, 240)}
FMT_DATE = {"bold": True, "font_size": 7, "color": RGBColor(0, 176, 240)} 
FMT_TEXT = {"font_size": 7}

def generate_replacements_slide_1(data):
    """
    Generates the replacements dict for Slide 1 (Overview).
    Maps Use Cases to {{SALES USE CASE Title X}}, {{SDX}}, {{SAX}} etc.
    """
    replacements = {}
    
    # LoB Key Mapping
    lob_map = {
        "Marketing": {"prefix": "Marketing", "date_prefix": "M"}, # {{Marketing USE CASE Title X}}
        "Sales": {"prefix": "SALES", "date_prefix": "S"},         # {{SALES USE CASE Title X}} note CAPS difference in template
        "Compliance": {"prefix": "Compliance", "date_prefix": "CO"},
        "Customer Success": {"prefix": "Customer Success", "date_prefix": "CU"},
    }
    
    for lob_name, config in lob_map.items():
        use_cases = data.get(lob_name, [])
        prefix = config["prefix"]
        d_pref = config["date_prefix"]
        
        for i, uc in enumerate(use_cases):
            num = i + 1
            # Title
            key_title = f"{{{{{prefix} USE CASE Title {num}}}}}"
            replacements[key_title] = {"text": uc.title, "formatting": FMT_TITLE}
            
            # Delivery Date
            key_del = f"{{{{{d_pref}D{num}}}}}"
            replacements[key_del] = {"text": uc.delivery_date, "formatting": FMT_DATE}
            
            # Adoption Date
            key_adopt = f"{{{{{d_pref}A{num}}}}}"
            replacements[key_adopt] = {"text": uc.adoption_date, "formatting": FMT_DATE}
            
    return replacements

def run():
    print("Loading data...")
    data = load_data(CSV_FILE)
    
    print("Opening PPTX...")
    prs = Presentation(INPUT_PPTX)
    
    # --- Process Slide 1 (Overview) ---
    print("Processing Slide 1...")
    rep_s1 = generate_replacements_slide_1(data)
    for shape in prs.slides[0].shapes:
        replace_text_in_shape(shape, rep_s1)
    
    # --- Process Slide 2 (Heatmap) ---
    print("Processing Slide 2...")
    slide2 = prs.slides[1]
    marketing_cases = data.get("Marketing", [])
    
    for shape in slide2.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    tf_text = cell.text_frame.text
                    
                    # Check which UC this cell belongs to
                    for i, uc in enumerate(marketing_cases):
                        num = i + 1
                        title_key = f"{{{{Marketing USE CASE Title {num}}}}}"
                        
                        if title_key in tf_text:
                            print(f"  Found Marketing UC {num} block: {uc.title}")
                            cell_replacements = {
                                title_key: {"text": uc.title, "formatting": FMT_TITLE},
                                f"{{{{StatusupdateUC{num}Marketing}}}}": {"text": uc.status_update, "formatting": FMT_TEXT, "is_bullet": True},
                                "{{UseCaseOwnerMarketing}}": {"text": uc.owner, "formatting": FMT_TEXT}
                            }
                            process_text_frame(cell.text_frame, cell_replacements)
                            break 

    # --- One Pagers ---
    print("Processing One Pagers...")
    # Find the Template Slide (Last Slide)
    num_slides = len(prs.slides)
    template_index = num_slides - 1 
    print(f"  Template slide index: {template_index}")
    
    # Iterate ALL use cases
    all_cases = []
    for cases in data.values():
        all_cases.extend(cases)
        
    for i, uc in enumerate(all_cases):
        print(f"  Generating One Pager for: {uc.title}")
        
        # Duplicate Template
        new_slide = duplicate_slide(prs, template_index)
        
        op_replacements = {
            "{{UseCaseOnePagerTitel1}}": {"text": uc.title, "formatting": FMT_TITLE},
            "{{UseCaseOnePagerPB1}}": {"text": uc.problem_statement, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerScope1}}": {"text": uc.scope, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerV&KPI1}}": {"text": uc.value_kpis, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerBU1}}": {"text": uc.business_unit, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerOwner1}}": {"text": uc.owner, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerScopeBC}}": {"text": uc.business_contacts, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerScopeAFK}}": {"text": uc.affected_key_users, "formatting": FMT_TEXT},
            "{{UseCaseOnePagerBSU1}}": {"text": "N/A", "formatting": FMT_TEXT} 
        }
        
        for shape in new_slide.shapes:
            replace_text_in_shape(shape, op_replacements)
    
    print(f"Saving to {OUTPUT_PPTX}...")
    prs.save(OUTPUT_PPTX)
    print("Done.")

if __name__ == "__main__":
    run()
